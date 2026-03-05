import os
import sys
import json
import re
import html
import base64
import requests
from pathlib import Path
from dotenv import load_dotenv

# Load .env file from the same directory as this script
load_dotenv(Path(__file__).parent.parent / ".env")

# Add parent directory to path for shared module
sys.path.insert(0, str(Path(__file__).parent.parent))
from sn_kb_shared import (
    upload_attachment,
    replace_base64_images,
    create_article,
    update_article,
)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff",
}


def get_config():
    instance = os.environ.get("SN_INSTANCE", "").rstrip("/")
    username = os.environ.get("SN_USERNAME", "")
    password = os.environ.get("SN_PASSWORD", "")
    kb_sys_id = os.environ.get("SN_KB_SYS_ID", "")

    missing = []
    if not instance:
        missing.append("SN_INSTANCE")
    if not username:
        missing.append("SN_USERNAME")
    if not password:
        missing.append("SN_PASSWORD")
    if not kb_sys_id:
        missing.append("SN_KB_SYS_ID")

    if missing:
        print(f"ERROR: Missing environment variables: {', '.join(missing)}")
        sys.exit(1)

    return instance, username, password, kb_sys_id


def rgb_to_hex(color_int):
    """Convert pymupdf integer color to hex string."""
    if color_int is None:
        return "#000000"
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    return f"#{r:02x}{g:02x}{b:02x}"


def extract_bg_regions(page, drawings=None):
    """Extract filled rectangles from PDF drawings to detect background colors and borders."""
    regions = []
    try:
        if drawings is None:
            drawings = page.get_drawings()
        for d in drawings:
            if d.get("fill") and d.get("rect"):
                rect = d["rect"]
                fill = d["fill"]
                # Convert fill color (tuple of 0-1 floats) to hex
                if isinstance(fill, (list, tuple)) and len(fill) >= 3:
                    r, g, b = int(fill[0] * 255), int(fill[1] * 255), int(fill[2] * 255)
                    color = f"#{r:02x}{g:02x}{b:02x}"
                else:
                    continue

                # Skip white/near-white backgrounds
                if r > 245 and g > 245 and b > 245:
                    continue

                stroke = d.get("color")
                border_color = None
                if stroke and isinstance(stroke, (list, tuple)) and len(stroke) >= 3:
                    sr, sg, sb = int(stroke[0] * 255), int(stroke[1] * 255), int(stroke[2] * 255)
                    if not (sr > 245 and sg > 245 and sb > 245):
                        border_color = f"#{sr:02x}{sg:02x}{sb:02x}"

                regions.append({
                    "bbox": (rect.x0, rect.y0, rect.x1, rect.y1),
                    "bg_color": color,
                    "border_color": border_color,
                    "width": rect.x1 - rect.x0,
                    "height": rect.y1 - rect.y0,
                })
    except Exception:
        pass
    return regions


def find_bg_for_element(elem_bbox, bg_regions):
    """Find the background region that contains this element."""
    ex0, ey0, ex1, ey1 = elem_bbox
    ecx = (ex0 + ex1) / 2
    ecy = (ey0 + ey1) / 2

    best = None
    best_area = float("inf")

    for region in bg_regions:
        rx0, ry0, rx1, ry1 = region["bbox"]
        # Check if element center is inside the region
        if rx0 <= ecx <= rx1 and ry0 <= ecy <= ry1:
            area = region["width"] * region["height"]
            # Prefer the smallest containing region (most specific)
            if area < best_area:
                best = region
                best_area = area

    return best


def detect_tables(page):
    """Detect grid-based tables on a PDF page using PyMuPDF's find_tables().

    Returns a list of dicts with keys: bbox, rows, col_count, row_count.
    Each row is a list of cell dicts: {text, bbox}.
    Filters out empty tables (fewer than 20% of cells have text).
    """
    tables = []
    try:
        found = page.find_tables()
        for table in found.tables:
            bbox = table.bbox  # already a tuple (x0, y0, x1, y1)
            extracted = table.extract()  # row-major: [[cell_text, ...], ...]
            row_count = table.row_count
            col_count = table.col_count

            # Build a cell bbox lookup from table.cells (column-major order)
            # cells order: col0_row0, col0_row1, ..., col1_row0, col1_row1, ...
            cell_bboxes = {}  # (row_idx, col_idx) -> (x0, y0, x1, y1)
            if hasattr(table, "cells") and table.cells:
                for idx, c in enumerate(table.cells):
                    if c:
                        col_idx = idx // row_count
                        row_idx = idx % row_count
                        cell_bboxes[(row_idx, col_idx)] = (c[0], c[1], c[2], c[3])

            rows = []
            total_cells = 0
            filled_cells = 0

            for row_idx, row_texts in enumerate(extracted):
                row_cells = []
                for col_idx, cell_text in enumerate(row_texts):
                    total_cells += 1
                    text = cell_text.strip() if cell_text else ""
                    if text:
                        filled_cells += 1
                    row_cells.append({
                        "text": text,
                        "bbox": cell_bboxes.get((row_idx, col_idx)),
                    })
                rows.append(row_cells)

            # Filter out mostly-empty tables
            if total_cells > 0 and (filled_cells / total_cells) < 0.2:
                continue

            tables.append({
                "bbox": bbox,
                "rows": rows,
                "col_count": col_count,
                "row_count": row_count,
            })
    except Exception:
        pass
    return tables


def is_inside_table(elem_bbox, tables):
    """Return True if the element's center point falls inside any detected table bbox."""
    ex0, ey0, ex1, ey1 = elem_bbox
    cx = (ex0 + ex1) / 2
    cy = (ey0 + ey1) / 2
    for t in tables:
        tx0, ty0, tx1, ty1 = t["bbox"]
        if tx0 <= cx <= tx1 and ty0 <= cy <= ty1:
            return True
    return False


def render_table_html(table_data, bg_regions):
    """Build an HTML <table> from detected table data with background colors."""
    parts = ['<table border="1" cellpadding="6" style="border-collapse:collapse;width:100%">']

    for row_idx, row in enumerate(table_data["rows"]):
        parts.append("<tr>")
        tag = "th" if row_idx == 0 else "td"

        for cell in row:
            if cell["bbox"] is None:
                # Merged/missing cell — skip
                continue
            text = html.escape(cell["text"]) if cell["text"] else "&nbsp;"
            # Check for background color on this cell
            bg = find_bg_for_element(cell["bbox"], bg_regions)
            style_parts = []
            if bg and bg.get("bg_color"):
                style_parts.append(f'background-color:{bg["bg_color"]}')
            if bg and bg.get("border_color"):
                style_parts.append(f'border-color:{bg["border_color"]}')
            style_attr = f' style="{";".join(style_parts)}"' if style_parts else ""
            parts.append(f"<{tag}{style_attr}>{text}</{tag}>")

        parts.append("</tr>")

    parts.append("</table>")
    return "\n".join(parts)


def detect_columns(blocks, page_width, bg_regions=None):
    """Detect if the page has a multi-column layout by analyzing text block positions and background regions."""
    if not blocks:
        return None

    # First check if there's a large background region that defines a sidebar
    if bg_regions:
        # Find the tallest narrow region (sidebar)
        sidebar_candidates = [
            r for r in bg_regions
            if r["height"] > 200 and 50 < r["width"] < page_width * 0.5
        ]
        if sidebar_candidates:
            sidebar = max(sidebar_candidates, key=lambda r: r["height"])
            return sidebar["bbox"][2]

    # Fallback: analyze text block positions
    text_blocks = [b for b in blocks if b["type"] == 0]
    if len(text_blocks) < 4:
        return None

    mid = page_width / 2
    left_blocks = [b for b in text_blocks if b["bbox"][0] < mid * 0.7]
    right_blocks = [b for b in text_blocks if b["bbox"][0] >= mid * 0.7]

    if len(left_blocks) >= 3 and len(right_blocks) >= 3:
        left_max_right = max(b["bbox"][2] for b in left_blocks)
        right_min_left = min(b["bbox"][0] for b in right_blocks)

        if right_min_left > left_max_right:
            return (left_max_right + right_min_left) / 2
        else:
            return mid * 0.7

    return None


def detect_text_align(elem_bbox, col_left, col_right):
    """Detect if a text block is centered within its column."""
    ex0, ex1 = elem_bbox[0], elem_bbox[2]
    text_width = ex1 - ex0
    col_width = col_right - col_left

    if col_width <= 0 or text_width <= 0:
        return "left"

    # If text takes up most of the column width, it's left-aligned (full paragraph)
    if text_width > col_width * 0.8:
        return "left"

    left_margin = ex0 - col_left
    right_margin = col_right - ex1

    # If both margins are substantial and roughly equal, it's centered
    if left_margin > 20 and right_margin > 20:
        margin_ratio = min(left_margin, right_margin) / max(left_margin, right_margin) if max(left_margin, right_margin) > 0 else 0
        if margin_ratio > 0.3:
            return "center"

    # If right margin is much larger, could be right-aligned
    if right_margin > left_margin * 3 and left_margin < 20:
        return "left"

    return "left"


def span_to_html(span):
    """Convert a pymupdf span to styled HTML."""
    text = span["text"]
    if not text.strip():
        return text

    styles = []
    font = span.get("font", "")
    size = span.get("size", 11)
    color = rgb_to_hex(span.get("color", 0))
    flags = span.get("flags", 0)

    is_bold = bool(flags & 2 ** 4) or "Bold" in font or "bold" in font.lower()
    is_italic = bool(flags & 2 ** 1) or "Italic" in font or "italic" in font.lower()

    if color and color != "#000000":
        styles.append(f"color:{color}")
    if size and size > 14:
        styles.append(f"font-size:{int(size)}px")

    style_attr = f' style="{";".join(styles)}"' if styles else ""

    if is_bold and is_italic:
        return f"<b><i><span{style_attr}>{text}</span></i></b>"
    elif is_bold:
        return f"<b><span{style_attr}>{text}</span></b>"
    elif is_italic:
        return f"<i><span{style_attr}>{text}</span></i>"
    elif styles:
        return f"<span{style_attr}>{text}</span>"
    else:
        return text


def block_to_html(block):
    """Convert a text block to HTML, preserving line breaks and span styling."""
    if block["type"] != 0:
        return ""

    lines_html = []
    for line in block.get("lines", []):
        spans_html = []
        for span in line.get("spans", []):
            spans_html.append(span_to_html(span))
        line_text = "".join(spans_html).strip()
        if line_text:
            lines_html.append(line_text)

    if not lines_html:
        return ""

    return "<br>".join(lines_html)


def extract_images_from_page(page, page_num, file_stem):
    """Extract images from a page with their bounding boxes and raw bytes."""
    images = []
    img_list = page.get_images(full=True)

    for img_idx, img_info in enumerate(img_list):
        xref = img_info[0]
        try:
            base_image = page.parent.extract_image(xref)
            if base_image:
                img_bytes = base_image["image"]
                ext = base_image.get("ext", "png")
                mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"

                # Find the image position on the page
                img_rects = page.get_image_rects(xref)
                if img_rects:
                    rect = img_rects[0]
                    bbox = (rect.x0, rect.y0, rect.x1, rect.y1)
                else:
                    bbox = None

                images.append({
                    "bytes": img_bytes,
                    "ext": ext,
                    "mime": mime,
                    "bbox": bbox,
                    "name": f"{file_stem}_p{page_num}_img{img_idx + 1}.{ext}",
                })
        except Exception:
            continue

    return images


def build_page_html(page, page_num, file_stem):
    """Build table-based HTML for a single PDF page."""
    page_dict = page.get_text("dict")
    blocks = page_dict.get("blocks", [])
    page_width = page_dict.get("width", 612)

    # Get drawings once and share between bg extraction and table detection
    drawings = page.get_drawings()

    # Extract background regions (colored rectangles, borders)
    bg_regions = extract_bg_regions(page, drawings=drawings)

    # Detect grid-based tables
    page_tables = detect_tables(page)

    # Extract images with positions
    page_images = extract_images_from_page(page, page_num, file_stem)

    # Build a unified list of elements (text blocks + images + tables) sorted by position
    elements = []

    for block in blocks:
        if block["type"] == 0:  # text block
            # Skip text blocks that fall inside a detected table
            if page_tables and is_inside_table(block["bbox"], page_tables):
                continue
            block_html = block_to_html(block)
            if block_html:
                bg = find_bg_for_element(block["bbox"], bg_regions)
                elements.append({
                    "type": "text",
                    "html": block_html,
                    "bbox": block["bbox"],
                    "bg": bg,
                })
        elif block["type"] == 1:  # image block embedded in page
            elements.append({
                "type": "image_placeholder",
                "bbox": block["bbox"],
                "bg": None,
            })

    for img in page_images:
        if img["bbox"]:
            elements.append({
                "type": "image",
                "data": img,
                "bbox": img["bbox"],
                "bg": None,
            })

    # Add detected tables as elements
    for tbl in page_tables:
        table_html = render_table_html(tbl, bg_regions)
        elements.append({
            "type": "table",
            "html": table_html,
            "bbox": tbl["bbox"],
            "bg": None,
        })

    # Detect column layout (uses bg_regions to find sidebar)
    split_x = detect_columns(blocks, page_width, bg_regions)

    if split_x:
        # Two-column layout using table
        left_elements = []
        right_elements = []
        full_width_elements = []

        for elem in elements:
            if elem["type"] == "table":
                # If the table spans most of the page width, render it full-width
                tbl_x0 = elem["bbox"][0]
                tbl_x1 = elem["bbox"][2]
                tbl_width = tbl_x1 - tbl_x0
                if tbl_width > page_width * 0.6:
                    full_width_elements.append(elem)
                    continue
            # Use the left edge (x0) for column assignment — more reliable than midpoint
            if elem["bbox"][0] < split_x:
                left_elements.append(elem)
            else:
                right_elements.append(elem)

        # Sort each column by vertical position
        left_elements.sort(key=lambda e: e["bbox"][1])
        right_elements.sort(key=lambda e: e["bbox"][1])
        full_width_elements.sort(key=lambda e: e["bbox"][1])

        left_width_pct = int((split_x / page_width) * 100)
        right_width_pct = 100 - left_width_pct

        # Determine sidebar background color from the largest bg region in the left column
        left_bg_style = ""
        sidebar_regions = [r for r in bg_regions
                          if r["bbox"][2] <= split_x + 10 and r["height"] > 100]
        if sidebar_regions:
            largest = max(sidebar_regions, key=lambda r: r["width"] * r["height"])
            left_bg_style = f' style="background-color:{largest["bg_color"]};'
            if largest.get("border_color"):
                left_bg_style += f' border:1px solid {largest["border_color"]};'
            left_bg_style += '"'

        left_html = render_column(left_elements, bg_regions, col_left=0, col_right=split_x)
        right_html = render_column(right_elements, bg_regions, col_left=split_x, col_right=page_width)

        column_html = (
            f'<table width="100%" cellpadding="8" cellspacing="0" border="0">'
            f'<tr>'
            f'<td width="{left_width_pct}%" valign="top"{left_bg_style}>{left_html}</td>'
            f'<td width="{right_width_pct}%" valign="top">{right_html}</td>'
            f'</tr></table>'
        )

        if full_width_elements:
            # Interleave full-width tables around the column layout based on position
            all_parts = []
            fw_before = [e for e in full_width_elements if e["bbox"][1] < left_elements[0]["bbox"][1]] if left_elements else full_width_elements
            fw_after = [e for e in full_width_elements if e not in fw_before]
            for elem in fw_before:
                all_parts.append(elem["html"])
            all_parts.append(column_html)
            for elem in fw_after:
                all_parts.append(elem["html"])
            return "\n".join(all_parts)

        return column_html
    else:
        # Single column layout
        elements.sort(key=lambda e: e["bbox"][1])
        return render_column(elements, bg_regions, col_left=0, col_right=page_width)


def render_column(elements, bg_regions=None, col_left=0, col_right=595):
    """Render a list of positioned elements as sequential HTML with background styling, spacing, and alignment."""
    html_parts = []
    prev_bottom = None

    for elem in elements:
        # Calculate vertical gap from previous element to preserve spacing
        margin_top = 0
        if prev_bottom is not None:
            gap = elem["bbox"][1] - prev_bottom
            if gap > 15:
                margin_top = int(gap * 0.8)
            elif gap > 8:
                margin_top = int(gap * 0.5)

        # Detect text alignment within column
        align = "left"
        if elem["type"] == "text":
            align = detect_text_align(elem["bbox"], col_left, col_right)

        style_parts = []
        if margin_top > 0:
            style_parts.append(f"margin-top:{margin_top}px")
        if align == "center":
            style_parts.append("text-align:center")

        if elem["type"] == "text":
            bg = elem.get("bg")
            if bg and bg.get("bg_color"):
                style_parts.append(f'background-color:{bg["bg_color"]}')
                style_parts.append("padding:4px 8px")
                if bg.get("border_color"):
                    style_parts.append(f'border:1px solid {bg["border_color"]}')
                style = ";".join(style_parts)
                html_parts.append(f'<div style="{style}">{elem["html"]}</div>')
            else:
                if style_parts:
                    style = ";".join(style_parts)
                    html_parts.append(f'<p style="{style}">{elem["html"]}</p>')
                else:
                    html_parts.append(f'<p>{elem["html"]}</p>')
        elif elem["type"] == "table":
            if margin_top > 0:
                html_parts.append(f'<div style="margin-top:{margin_top}px">{elem["html"]}</div>')
            else:
                html_parts.append(elem["html"])
        elif elem["type"] == "image":
            img_data = elem["data"]
            b64 = base64.b64encode(img_data["bytes"]).decode("ascii")
            width = int(elem["bbox"][2] - elem["bbox"][0])
            img_style_parts = ["max-width:100%"]
            if margin_top > 0:
                img_style_parts.append(f"margin-top:{margin_top}px")
            p_align = ' style="text-align:center"' if align == "center" else ""
            html_parts.append(
                f'<p{p_align}><img src="data:{img_data["mime"]};base64,{b64}" '
                f'width="{width}" style="{";".join(img_style_parts)}" /></p>'
            )

        prev_bottom = elem["bbox"][3]

    return "\n".join(html_parts)


def extract_pdf(file_path):
    import fitz  # pymupdf

    doc = fitz.open(file_path)
    file_stem = file_path.stem
    html_parts = []

    for page_num, page in enumerate(doc, 1):
        page_html = build_page_html(page, page_num, file_stem)
        if page_html:
            html_parts.append(page_html)
            if page_num < len(doc):
                html_parts.append('<hr style="border:1px solid #ccc; margin:20px 0;">')

    doc.close()
    return "\n".join(html_parts)


_DOCX_MAX_IMG_PX = 900    # max width or height for embedded images
_DOCX_IMG_QUALITY = 72    # JPEG quality for recompressed images


def _compress_docx_image(image):
    """Compress and resize an image extracted from a DOCX file."""
    import io
    import base64
    try:
        from PIL import Image as PILImage
        with image.open() as f:
            raw = f.read()
        pil = PILImage.open(io.BytesIO(raw))
        # Convert palette/RGBA to RGB for JPEG output
        if pil.mode in ("P", "RGBA", "LA"):
            pil = pil.convert("RGBA")
            bg = PILImage.new("RGB", pil.size, (255, 255, 255))
            bg.paste(pil, mask=pil.split()[3])
            pil = bg
        elif pil.mode != "RGB":
            pil = pil.convert("RGB")
        # Resize if too large
        w, h = pil.size
        if max(w, h) > _DOCX_MAX_IMG_PX:
            scale = _DOCX_MAX_IMG_PX / max(w, h)
            pil = pil.resize((int(w * scale), int(h * scale)), PILImage.LANCZOS)
        buf = io.BytesIO()
        pil.save(buf, format="JPEG", quality=_DOCX_IMG_QUALITY, optimize=True)
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        return {"src": f"data:image/jpeg;base64,{b64}"}
    except Exception:
        # Fall back to original if compression fails
        with image.open() as f:
            raw = f.read()
        import base64
        b64 = base64.b64encode(raw).decode("ascii")
        ct = getattr(image, "content_type", "image/png")
        return {"src": f"data:{ct};base64,{b64}"}


def extract_docx(file_path):
    import mammoth
    import re

    # Map custom/French DOCX styles to proper HTML elements
    style_map = """
p[style-name='Style2'] => h1:fresh
p[style-name='Style1'] => h2:fresh
p[style-name='TOC Heading'] => h2:fresh
p[style-name='header'] => h2:fresh
p[style-name='Header'] => h2:fresh
p[style-name='toc 1'] => p:fresh
p[style-name='TEXTE1'] => p:fresh
p[style-name='List Paragraph'] => ul > li:fresh
p[style-name='ListParagraph'] => ul > li:fresh
p[style-name='Paragraphedeliste'] => ul > li:fresh
"""

    convert_image = mammoth.images.img_element(_compress_docx_image)

    with open(file_path, "rb") as f:
        result = mammoth.convert_to_html(f, style_map=style_map, convert_image=convert_image)

    html = result.value

    # Replace non-breaking spaces with regular spaces
    html = html.replace("\u00a0", " ")

    # Remove empty paragraphs (including those containing only anchors or whitespace)
    html = re.sub(r"<p>\s*(<a[^>]+></a>\s*)*</p>", "", html)

    # Strip orphan bookmark anchors (no href, no text — Word cross-reference bookmarks)
    html = re.sub(r'<a\s+id="[^"]*"\s*>\s*</a>', "", html)

    # Collapse multiple consecutive spaces
    html = re.sub(r" {2,}", " ", html)

    # Improve table styling, promote header cells, add alternating row colors
    html = _style_docx_tables(html)

    return html


_TABLE_STYLE = (
    'border="1" cellpadding="8" cellspacing="0" '
    'style="border-collapse:collapse;width:100%;margin-bottom:16px;"'
)
_TH_STYLE = 'style="background-color:#d9e1f2;font-weight:bold;text-align:left;padding:6px 8px;"'
_TD_STYLE = 'style="padding:6px 8px;vertical-align:top;"'


def _style_docx_tables(html):
    """Add borders/styling to tables, promote header cells, add alternating row colors."""
    import re

    def process_table(m):
        table_html = m.group(0)
        # Add styling attributes to the <table> tag
        table_html = re.sub(
            r"<table(?:\s[^>]*)?>",
            f"<table {_TABLE_STYLE}>",
            table_html,
            count=1,
        )
        # Collect all <tr>...</tr> blocks
        rows = list(re.finditer(r"<tr>(.*?)</tr>", table_html, re.DOTALL))
        has_header = False

        if rows:
            first_row = rows[0].group(0)
            cells = re.findall(r"<td[^>]*>(.*?)</td>", first_row, re.DOTALL)
            # Promote bold-only first row to <th>
            if cells and all(
                re.match(r"\s*<p>(\s*<a[^>]+></a>\s*)*\s*<strong>.*?</strong>\s*</p>\s*$", c, re.DOTALL)
                for c in cells
            ):
                new_row = re.sub(
                    r"<td[^>]*>\s*<p>(?:\s*<a[^>]+></a>\s*)*\s*<strong>(.*?)</strong>\s*</p>\s*</td>",
                    lambda mm: f'<th {_TH_STYLE}>{mm.group(1)}</th>',
                    first_row,
                    flags=re.DOTALL,
                )
                table_html = table_html.replace(first_row, new_row, 1)
                has_header = True

        # Add alternating row background to data rows (skip header row)
        if has_header and len(rows) > 2:
            data_rows = list(re.finditer(r"<tr>(.*?)</tr>", table_html, re.DOTALL))
            for idx, row_m in enumerate(data_rows[1:], start=0):
                bg = "#f2f2f2" if idx % 2 == 0 else "#ffffff"
                styled_tr = f'<tr style="background-color:{bg};">'
                table_html = table_html.replace(row_m.group(0),
                                                row_m.group(0).replace("<tr>", styled_tr, 1), 1)

        # Add vertical-align:top style to all remaining <td> tags
        table_html = re.sub(r"<td(?:\s[^>]*)?>", f"<td {_TD_STYLE}>", table_html)
        return table_html

    return re.sub(r"<table.*?</table>", process_table, html, flags=re.DOTALL)


def extract_xlsx(file_path):
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    html_parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        html_parts.append(f"<h2>Sheet: {sheet}</h2>")
        html_parts.append("<table border='1' cellpadding='4' cellspacing='0'>")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                row_html = "".join(f"<td>{c}</td>" for c in cells)
                html_parts.append(f"<tr>{row_html}</tr>")
        html_parts.append("</table>")
    wb.close()
    return "\n".join(html_parts)


def extract_pptx(file_path):
    from pptx import Presentation

    prs = Presentation(file_path)
    html_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_texts.append(f"<p>{paragraph.text}</p>")
        if slide_texts:
            html_parts.append(f"<h2>Slide {i}</h2>")
            html_parts.extend(slide_texts)
    return "\n".join(html_parts)


def extract_image_file(file_path):
    """For standalone image files, encode as base64 HTML img."""
    ext_to_mime = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".bmp": "image/bmp",
        ".tiff": "image/tiff",
    }
    mime = ext_to_mime.get(file_path.suffix.lower(), "image/png")
    img_data = file_path.read_bytes()
    b64 = base64.b64encode(img_data).decode("ascii")
    return f'<p><img src="data:{mime};base64,{b64}" style="max-width:100%;" /></p>'


def extract_html(file_path):
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".xlsx":
        return extract_xlsx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}:
        return extract_image_file(file_path)
    else:
        return None


def main():
    if len(sys.argv) < 2:
        print("Usage: python upload_to_kb.py <folder_path>")
        sys.exit(1)

    folder_path = Path(sys.argv[1])
    if not folder_path.is_dir():
        print(f"ERROR: '{folder_path}' is not a valid directory.")
        sys.exit(1)

    instance, username, password, kb_sys_id = get_config()

    files = [f for f in folder_path.iterdir() if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]

    if not files:
        print(f"No supported files found in '{folder_path}'.")
        print(f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        sys.exit(0)

    print(f"Found {len(files)} file(s) to process.\n")

    results = {"success": [], "failed": []}

    for file_path in sorted(files):
        print(f"Processing: {file_path.name}...")

        try:
            html = extract_html(file_path)
            if not html or not html.strip():
                print("  SKIPPED (no content extracted)")
                results["failed"].append((file_path.name, "No content could be extracted"))
                continue

            title = file_path.stem.replace("_", " ").replace("-", " ").title()

            # Step 1: Create the article with a placeholder
            print("  Creating KB article...", end=" ")
            article_sys_id, number = create_article(
                instance, username, password, kb_sys_id, title, "<p>Uploading content...</p>"
            )

            if not article_sys_id:
                print(f"FAILED ({number[:100]})")
                results["failed"].append((file_path.name, number[:200]))
                continue

            print(f"OK ({number})")

            # Step 2: Upload embedded images as attachments and replace base64 src
            has_images = "data:image/" in html
            if has_images:
                print("  Uploading embedded images...", end=" ")
                html, img_count = replace_base64_images(
                    html, instance, username, password, article_sys_id, file_path.stem
                )
                print(f"{img_count} image(s) uploaded")

            # Step 3: Update the article with the final HTML
            print("  Updating article content...", end=" ")
            if update_article(instance, username, password, article_sys_id, html):
                print("OK")
                results["success"].append((file_path.name, number))
            else:
                print("FAILED (could not update article body)")
                results["failed"].append((file_path.name, "Article created but body update failed"))

        except Exception as e:
            print(f"  ERROR ({e})")
            results["failed"].append((file_path.name, str(e)))

        print()

    print(f"--- Summary ---")
    print(f"Successful: {len(results['success'])}")
    print(f"Failed:     {len(results['failed'])}")

    if results["success"]:
        print("\nUploaded articles:")
        for name, number in results["success"]:
            print(f"  - {name} -> {number}")

    if results["failed"]:
        print("\nFailed files:")
        for name, reason in results["failed"]:
            print(f"  - {name}: {reason}")


if __name__ == "__main__":
    main()
